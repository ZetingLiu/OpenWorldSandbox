"""将 OpenWorldSandbox HTML 计划稿重建为原生可编辑 PPTX。"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "openworldsandbox_full_plan_slides.pptx"

SW, SH = 13.333, 7.5
BG = "FFFFFF"
CARD = "F7F9FC"
CARD_2 = "EDF3F8"
INK = "172033"
MUTED = "46566E"
FAINT = "74839A"
LINE = "D8E1EC"
CYAN = "34D3EE"
TEAL = "34D399"
VIOLET = "A78BFA"
AMBER = "FBBF24"
RED = "FB7185"
FONT = "Microsoft YaHei"
MONO = "Consolas"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=14,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.03,
    name=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        box.name = name
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = 1.1
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, *, size=14, align=PP_ALIGN.LEFT, name=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        box.name = name
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, color, bold in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_shape(
    slide,
    shape_type,
    x,
    y,
    w,
    h,
    *,
    fill=CARD,
    line=LINE,
    line_width=1.0,
    radius_name=None,
):
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if radius_name:
        shape.name = radius_name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_card(slide, x, y, w, h, *, fill=CARD, line=LINE, name=None):
    return add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=fill,
        line=line,
        radius_name=name,
    )


def add_line(
    slide,
    x1,
    y1,
    x2,
    y2,
    *,
    color=FAINT,
    width=1.5,
    dash=False,
    name=None,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    if name:
        line.name = name
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_arrow(slide, x1, y1, x2, y2, *, color=FAINT, width=1.5, name=None):
    add_line(slide, x1, y1, x2, y2, color=color, width=width, name=name)
    angle = 0 if x2 >= x1 else 180
    tri = add_shape(
        slide,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        x2 - 0.08,
        y2 - 0.08,
        0.16,
        0.16,
        fill=color,
        line=color,
    )
    tri.rotation = 90 + angle
    return tri


def add_badge(slide, text, x, y, w, *, color=CYAN, name=None):
    add_card(slide, x, y, w, 0.34, fill=CARD_2, line=color, name=name)
    return add_text(
        slide,
        text,
        x + 0.06,
        y + 0.01,
        w - 0.12,
        0.30,
        size=9.5,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_background(slide, page_no):
    bg = add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, fill=BG, line=BG)
    bg.name = "背景"
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    add_text(
        slide,
        f"{page_no:02d}",
        12.35,
        7.08,
        0.45,
        0.18,
        size=8.5,
        color=FAINT,
        align=PP_ALIGN.RIGHT,
        name="页码",
    )


def add_header(slide, page_no, kicker, title, badge=None, badge_color=CYAN):
    add_background(slide, page_no)
    add_text(slide, kicker, 0.68, 0.22, 5.0, 0.23, size=9, color=FAINT, bold=True)
    add_text(slide, title, 0.68, 0.47, 10.8, 0.45, size=21.5, bold=True, name="标题")
    if badge:
        width = max(1.35, min(2.75, 0.105 * len(badge) + 0.55))
        add_badge(slide, badge, 12.65 - width, 0.44, width, color=badge_color)
    add_line(slide, 0.68, 1.05, 12.65, 1.05, color=LINE, width=1)


def add_bullet(slide, text, x, y, w, *, color=CYAN, size=12.5, symbol="✓"):
    add_text(slide, symbol, x, y, 0.28, 0.28, size=size, color=color, bold=True)
    add_text(slide, text, x + 0.34, y, w - 0.34, 0.36, size=size, color=MUTED)


def add_code(slide, text, x, y, w, h, *, name=None, size=10.5):
    add_card(slide, x, y, w, h, fill="F4F7FA", line=LINE, name=name)
    return add_text(
        slide,
        text,
        x + 0.18,
        y + 0.12,
        w - 0.36,
        h - 0.24,
        size=size,
        color=MUTED,
        font=MONO,
        valign=MSO_ANCHOR.TOP,
        margin=0,
    )


def add_table(slide, rows, cols, x, y, w, h, data, widths=None, font_size=10.5):
    shape = slide.shapes.add_table(
        rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.name = "可编辑数据表"
    table = shape.table
    if widths:
        for idx, width in enumerate(widths):
            table.columns[idx].width = Inches(width)
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(CARD_2 if row_idx == 0 else CARD)
            cell.border_left = None
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = paragraph.runs[0]
            run.font.name = FONT
            run.font.size = Pt(font_size if row_idx else font_size - 0.5)
            run.font.bold = row_idx == 0 or col_idx == 0
            run.font.color.rgb = rgb(INK if row_idx == 0 or col_idx == 0 else MUTED)
    return shape


def add_loop_node(slide, title, subtitle, x, y, d, color):
    add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill=CARD_2, line=color, line_width=1.7)
    add_text(
        slide,
        title,
        x + 0.12,
        y + d * 0.34,
        d - 0.24,
        0.30,
        size=15,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        subtitle,
        x + 0.14,
        y + d * 0.55,
        d - 0.28,
        0.32,
        size=9.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def add_summary_bar(slide, text, y, *, color=CYAN):
    add_card(slide, 0.95, y, 11.43, 0.62, fill=CARD, line=LINE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.95, y, 0.08, 0.62, fill=color, line=color)
    add_text(slide, text, 1.18, y + 0.08, 10.95, 0.44, size=12.5, color=MUTED)


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    # 01 封面
    slide = prs.slides.add_slide(blank)
    add_background(slide, 1)
    add_text(slide, "转正答辩 · 项目计划", 0.75, 0.55, 4.8, 0.28, size=11, color=CYAN, bold=True)
    add_text(slide, "OpenWorldSandbox", 0.75, 1.15, 10.7, 0.72, size=38, bold=True, name="主标题")
    add_text(slide, "为具身大脑构建", 0.75, 1.93, 10.7, 0.62, size=31, color=CYAN, bold=True)
    add_text(slide, "可交互、可验证的世界", 0.75, 2.52, 11.1, 0.67, size=34, color=VIOLET, bold=True)
    add_text(
        slide,
        "从多轮具身推理评测，到 Pelican-VL 的 Agentic RL 执行环境。",
        0.78,
        3.42,
        10.8,
        0.42,
        size=16,
        color=MUTED,
    )
    route = ["为什么需要", "项目定位", "系统工作流", "真实 Case", "Pelican-VL", "建设计划"]
    x = 0.78
    for idx, item in enumerate(route):
        width = 1.42 if idx != 2 else 1.58
        add_card(slide, x, 4.35, width, 0.5, fill=CARD, line=LINE)
        add_text(slide, item, x + 0.04, 4.43, width - 0.08, 0.32, size=10.5, align=PP_ALIGN.CENTER)
        x += width
        if idx < len(route) - 1:
            add_text(slide, "→", x + 0.03, 4.44, 0.32, 0.3, size=13, color=FAINT, align=PP_ALIGN.CENTER)
            x += 0.38
    add_line(slide, 0.78, 5.55, 12.2, 5.55, color=LINE)
    add_text(
        slide,
        "OpenWorldSandbox · 2026.08 · 项目计划版",
        0.78,
        6.55,
        5.2,
        0.24,
        size=9,
        color=FAINT,
    )

    # 02 为什么需要
    slide = prs.slides.add_slide(blank)
    add_header(
        slide,
        2,
        "01 · 为什么需要",
        "Brain Model 的能力，必须在多轮交互中才能被测试",
        "Static Answer ≠ Embodied Reasoning",
        CYAN,
    )
    add_text(slide, "静态测试：只看到一次回答", 1.17, 1.35, 4.7, 0.3, size=12, color=FAINT, bold=True, align=PP_ALIGN.CENTER)
    for x, title, sub in [
        (0.85, "固定输入", "一张图 / 一个问题"),
        (2.75, "一次回答", "世界不会变化"),
        (4.65, "对 / 错", "无法看到过程"),
    ]:
        add_card(slide, x, 2.1, 1.45, 0.88)
        add_text(slide, title, x + 0.08, 2.28, 1.29, 0.27, size=12, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x + 0.08, 2.57, 1.29, 0.22, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.30, 2.54, 2.65, 2.54, color=FAINT)
    add_arrow(slide, 4.20, 2.54, 4.55, 2.54, color=FAINT)
    add_text(slide, "模型知道“应该怎么做”", 1.35, 3.55, 4.2, 0.3, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "不代表它能连续把任务做完", 1.35, 3.88, 4.2, 0.25, size=10.5, color=FAINT, align=PP_ALIGN.CENTER)
    add_line(slide, 6.55, 1.38, 6.55, 5.72, color=LINE)
    add_text(slide, "具身 Sandbox：观察—行动—反馈循环", 7.16, 1.35, 5.1, 0.3, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_loop_node(slide, "Brain Model", "理解 · 规划 · 决策", 7.08, 2.08, 1.82, TEAL)
    add_loop_node(slide, "Sandbox", "执行 · 更新 · 反馈", 10.42, 2.08, 1.82, VIOLET)
    add_arrow(slide, 8.78, 2.35, 10.34, 2.35, color=CYAN, width=2.2)
    add_text(slide, "高层动作", 9.15, 2.02, 0.85, 0.24, size=10, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 10.47, 3.76, 8.86, 3.76, color=VIOLET, width=2.2)
    add_text(slide, "新观察 / 失败原因", 9.18, 3.88, 1.35, 0.24, size=10, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 7.20, 4.62, 4.92, 0.88)
    add_text(slide, "只有在循环中，能力缺口才会暴露", 7.42, 4.78, 4.48, 0.26, size=12, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "状态跟踪 · 多步规划 · 条件判断 · 失败后重规划", 7.42, 5.10, 4.48, 0.23, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_summary_bar(
        slide,
        "核心结论：具身推理发生在“观察—行动—反馈”循环中，Brain Model 需要一个会随动作改变的符号级环境。",
        6.18,
    )

    # 03 项目定位
    slide = prs.slides.add_slide(blank)
    add_header(slide, 3, "02 · 项目定位", "OpenWorldSandbox 补上多轮推理与高层行动这一层", "各自解决不同问题", TEAL)
    add_card(slide, 4.83, 1.30, 3.68, 0.75, fill=CARD_2, line=CYAN)
    add_text(slide, "Benchmark", 4.98, 1.44, 3.38, 0.25, size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "定义考什么、任务集合和什么算成功", 4.98, 1.72, 3.38, 0.21, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 6.67, 2.05, 6.67, 2.50, color=CYAN, width=2)
    add_loop_node(slide, "Brain Model", "观察场景 · 决定下一步", 1.15, 2.48, 2.20, TEAL)
    add_card(slide, 4.65, 2.50, 4.05, 2.18, fill=CARD_2, line=VIOLET)
    add_text(slide, "OpenWorldSandbox", 4.90, 2.82, 3.55, 0.38, size=18, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "符号级交互环境", 4.90, 3.21, 3.55, 0.28, size=12.5, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "保存世界状态 · 执行高层动作", 4.90, 3.62, 3.55, 0.26, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "返回新观察与失败原因", 4.90, 3.92, 3.55, 0.26, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "测试：下一步决定做什么", 4.90, 4.25, 3.55, 0.25, size=10.5, color=AMBER, align=PP_ALIGN.CENTER)
    add_card(slide, 10.05, 2.50, 2.45, 2.18)
    add_text(slide, "物理仿真引擎", 10.23, 2.82, 2.09, 0.36, size=16, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "关节 · 碰撞 · 摩擦", 10.23, 3.30, 2.09, 0.24, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "抓取 · 运动轨迹 · 画面", 10.23, 3.62, 2.09, 0.24, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "测试：身体如何完成动作", 10.23, 4.12, 2.09, 0.26, size=10, color=AMBER, align=PP_ALIGN.CENTER)
    add_arrow(slide, 3.25, 3.03, 4.57, 3.03, color=CYAN, width=2)
    add_text(slide, "高层动作", 3.57, 2.73, 0.82, 0.23, size=9.5, color=CYAN, align=PP_ALIGN.CENTER)
    add_arrow(slide, 4.57, 4.13, 3.25, 4.13, color=VIOLET, width=2)
    add_text(slide, "反馈", 3.70, 4.28, 0.5, 0.2, size=9.5, color=VIOLET, align=PP_ALIGN.CENTER)
    add_arrow(slide, 8.70, 3.60, 9.96, 3.60, color=FAINT, width=1.5)
    add_text(slide, "能力下沉", 9.02, 3.30, 0.75, 0.2, size=9, color=FAINT, align=PP_ALIGN.CENTER)
    add_summary_bar(slide, "Benchmark 是试卷 · OpenWorldSandbox 是可交互的考场 · 仿真引擎验证身体和物理。三者不是替代关系。", 5.75, color=VIOLET)

    # 04 工作流
    slide = prs.slides.add_slide(blank)
    add_header(slide, 4, "03 · 系统架构", "从任务定义到能力报告：项目覆盖完整工作流", "Prepare → Interact → Evaluate", VIOLET)
    stages = [
        ("场景与任务 JSON", "房间 / 物体 / 初始状态\n指令 / 目标 / 参考动作", CYAN),
        ("ows compile", "检查数据是否合法\n回放参考动作验证有解", CYAN),
        ("SQLite 初始世界", "固定起点，可以重复运行\nowb env start", TEAL),
        ("Brain Model 运行", "观察 → 动作 → 反馈\n直到完成或超过步数", TEAL),
        ("验证与报告", "检查最终世界是否满足目标\n分析失败、效率和能力维度", VIOLET),
    ]
    xs = [0.45, 2.98, 5.51, 8.04, 10.57]
    for idx, ((title, desc, color), x) in enumerate(zip(stages, xs)):
        add_card(slide, x, 1.72, 2.25, 1.28, fill=CARD_2, line=color)
        add_text(slide, title, x + 0.12, 1.95, 2.01, 0.30, size=12.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.12, 2.34, 2.01, 0.48, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
        if idx < 4:
            add_arrow(slide, x + 2.25, 2.36, xs[idx + 1] - 0.08, 2.36, color=FAINT)
    add_text(slide, "离线准备", 1.58, 1.28, 2.25, 0.23, size=10.5, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "在线交互", 6.78, 1.28, 2.25, 0.23, size=10.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "离线评测", 10.57, 1.28, 2.25, 0.23, size=10.5, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 3.10, 3.75, 7.15, 1.25)
    add_text(slide, "运行产物", 3.32, 3.95, 6.71, 0.25, size=13, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "trajectory.json：每一步做了什么、为什么、环境如何反馈", 3.32, 4.28, 6.71, 0.25, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "final.db：任务结束后，世界中每个物体和设备的真实状态", 3.32, 4.58, 6.71, 0.25, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 9.16, 3.00, 8.84, 3.68, color=FAINT)
    add_arrow(slide, 10.25, 4.37, 11.45, 3.08, color=FAINT)
    add_line(slide, 11.69, 3.00, 11.69, 5.77, color=AMBER, width=1.8, dash=True)
    add_line(slide, 11.69, 5.77, 1.25, 5.77, color=AMBER, width=1.8, dash=True)
    add_arrow(slide, 1.25, 5.77, 1.25, 3.08, color=AMBER, width=1.8)
    add_summary_bar(slide, "评测结果回到任务设计与模型研究，形成下一轮迭代。", 6.14, color=AMBER)

    # 05 运行时架构
    slide = prs.slides.add_slide(blank)
    add_header(slide, 5, "04 · 运行时架构", "模型不能直接改世界：动作要经过规则与常识检查", "可控 · 可追踪 · 可验证", AMBER)
    add_loop_node(slide, "Brain Model", "选择高层动作", 0.55, 1.78, 1.65, TEAL)
    add_card(slide, 3.00, 1.90, 2.03, 1.38, fill=CARD_2, line=CYAN)
    add_text(slide, "HTTP / MCP Server", 3.15, 2.14, 1.73, 0.30, size=12.5, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "暴露 17 个动作\n统一参数与返回格式", 3.15, 2.56, 1.73, 0.44, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, 5.80, 1.22, 2.73, 0.62, fill=CARD_2, line=AMBER)
    add_text(slide, "Commonsense Judge LLM（可选）", 5.96, 1.34, 2.41, 0.23, size=10.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "提交前判断物理动作是否合理", 5.96, 1.57, 2.41, 0.18, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, 6.15, 2.10, 2.03, 1.38, fill=CARD_2, line=VIOLET)
    add_text(slide, "Action Rules", 6.30, 2.33, 1.73, 0.30, size=13.5, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "检查位置、双手与容器状态\n成功才允许改变世界", 6.30, 2.75, 1.73, 0.44, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.OVAL, 9.10, 1.73, 2.35, 1.82, fill=CARD_2, line=AMBER, line_width=1.7)
    add_text(slide, "SQLite World State", 9.31, 2.15, 1.93, 0.30, size=13.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "房间 · 物体 · 设备 · 双手\n世界事实只保存在这里", 9.31, 2.58, 1.93, 0.47, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, 11.77, 3.86, 1.05, 0.70)
    add_text(slide, "Evaluator", 11.87, 4.00, 0.85, 0.22, size=10.5, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "查终态", 11.87, 4.25, 0.85, 0.18, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.20, 2.35, 2.92, 2.35, color=CYAN, width=2)
    add_arrow(slide, 5.03, 2.70, 6.07, 2.70, color=FAINT)
    add_arrow(slide, 8.18, 2.70, 9.02, 2.70, color=CYAN, width=2)
    add_arrow(slide, 4.68, 1.90, 5.72, 1.53, color=AMBER)
    add_arrow(slide, 7.16, 1.84, 7.16, 2.02, color=AMBER)
    add_arrow(slide, 11.21, 3.32, 11.75, 4.05, color=FAINT)
    add_arrow(slide, 3.05, 3.08, 2.05, 3.60, color=VIOLET, width=2)
    rows = [
        ("世界怎么改变", "常识判断与固定规则均通过才更新；任一环节拒绝，数据库不变。", CYAN),
        ("模型能看到什么", "观察由数据库生成，只展示当前位置和打开容器中的物体。", TEAL),
        ("任务怎么算完成", "评测器查询 final.db，而不是相信模型说“我完成了”。", VIOLET),
    ]
    y = 4.92
    for title, desc, color in rows:
        add_card(slide, 0.85, y, 11.65, 0.50, fill=CARD, line=LINE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.85, y, 0.07, 0.50, fill=color, line=color)
        add_text(slide, title, 1.08, y + 0.08, 1.45, 0.30, size=11.5, color=color, bold=True)
        add_text(slide, desc, 2.65, y + 0.08, 9.55, 0.30, size=10.5, color=MUTED)
        y += 0.58

    # 06 当前进展
    slide = prs.slides.add_slide(blank)
    add_header(slide, 6, "05 · 当前进展", "已经跑通最小闭环，但距离可持续使用还有建设空间", "当前仓库真实状态", TEAL)
    for x, number, label, color in [
        (1.55, "2", "家庭 / 商超场景", CYAN),
        (5.55, "7", "可编译任务", TEAL),
        (9.55, "17", "高层语义动作", VIOLET),
    ]:
        add_shape(slide, MSO_SHAPE.OVAL, x, 1.35, 1.35, 1.35, fill=CARD_2, line=color, line_width=1.7)
        add_text(slide, number, x + 0.15, 1.58, 1.05, 0.44, size=25, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.05, 2.12, 1.25, 0.28, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_line(slide, 2.90, 2.02, 5.55, 2.02, color=LINE, width=1.5)
    add_line(slide, 6.90, 2.02, 9.55, 2.02, color=LINE, width=1.5)
    for x, title, items, color, symbol in [
        (0.85, "已经具备", [
            "JSON → SQLite 的编译与可解性检查",
            "HTTP / MCP 多轮交互与可选常识检查",
            "完整轨迹、最终状态和程序化目标判定",
            "成功率、路径效率和失败诊断报告",
        ], TEAL, "✓"),
        (6.78, "仍需建设", [
            "任务数量、难度层次和扰动覆盖",
            "会随时间或状态变化的环境事件",
            "核心规则回归测试与标准实验协议",
            "reset / step、奖励和并行环境能力",
        ], AMBER, "+"),
    ]:
        add_card(slide, x, 3.05, 5.70, 2.48)
        add_text(slide, title, x + 0.25, 3.25, 2.0, 0.28, size=13.5, bold=True)
        for idx, item in enumerate(items):
            add_bullet(slide, item, x + 0.28, 3.70 + idx * 0.42, 5.12, color=color, size=10.5, symbol=symbol)
    add_summary_bar(slide, "当前判断：主流程已经可用，下一步重点是把“能跑”变成“稳定、可比较、可扩展”。", 6.15, color=TEAL)

    # 07 Case 任务 JSON
    slide = prs.slides.add_slide(blank)
    add_header(slide, 7, "CASE · 01 / 04", "任务从两个 JSON 开始：世界是什么，目标是什么", "真实项目文件", CYAN)
    add_text(slide, "场景 JSON：定义初始世界", 0.85, 1.36, 5.65, 0.32, size=13, bold=True)
    add_code(
        slide,
        '{\n  "id": "clothes_01",\n  "class": "clothing",\n  "name": "衬衫",\n  "on": "bed_01",\n  "states": {"cleanliness": "dirty"}\n}\n{\n  "id": "washing_machine_01",\n  "open_state": "closed",\n  "device_state": "off",\n  "properties": ["can_contain", "can_wash"]\n}',
        0.85,
        1.78,
        5.72,
        4.65,
        name="场景 JSON 代码",
        size=10.8,
    )
    add_text(slide, "来源：data/scenarios/home_01.json", 1.03, 6.47, 4.5, 0.20, size=8.5, color=FAINT)
    add_text(slide, "任务 JSON：定义完成条件", 6.78, 1.36, 5.65, 0.32, size=13, bold=True)
    add_code(
        slide,
        '"instruction": "把两件脏衣物放入\\n洗衣机并启动洗涤",\n\n"goal": {"all_of": [\n  {"entity": "clothes_01",\n   "field": "container_id",\n   "value": "washing_machine_01"},\n  {"entity": "clothes_02",\n   "field": "container_id",\n   "value": "washing_machine_01"},\n  {"entity": "washing_machine_01",\n   "field": "device_state",\n   "value": "running"}\n]}',
        6.78,
        1.78,
        5.72,
        4.65,
        name="任务 JSON 代码",
        size=10.8,
    )
    add_text(slide, "来源：data/tasks/home/home_01_laundry_basic.json", 6.96, 6.47, 5.0, 0.20, size=8.5, color=FAINT)

    # 08 Case Compile
    slide = prs.slides.add_slide(blank)
    add_header(slide, 8, "CASE · 02 / 04", "Compile 把 JSON 变成可执行的 SQLite 世界", "solvable = true", TEAL)
    add_code(
        slide,
        "SELECT id, area_id, container_id, on_id, open_state, device_state\nFROM entities\nWHERE id IN ('clothes_01', 'clothes_02', 'washing_machine_01');",
        0.85,
        1.35,
        11.65,
        0.88,
        name="SQL 查询",
        size=9.5,
    )
    table_data = [
        ["id", "area_id", "container_id", "on_id", "open_state", "device_state"],
        ["clothes_01", "bedroom", "NULL", "bed_01", "NULL", "NULL"],
        ["clothes_02", "bedroom", "NULL", "bed_01", "NULL", "NULL"],
        ["washing_machine_01", "balcony", "NULL", "NULL", "closed", "off"],
    ]
    add_table(slide, 4, 6, 0.85, 2.48, 11.65, 1.63, table_data, widths=[2.25, 1.55, 2.35, 1.75, 1.85, 1.90], font_size=10.2)
    add_card(slide, 0.85, 4.55, 11.65, 1.22)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.85, 4.55, 0.08, 1.22, fill=CYAN, line=CYAN)
    add_text(slide, "Compile 做了什么", 1.18, 4.79, 2.30, 0.34, size=14, color=CYAN, bold=True)
    add_text(
        slide,
        "检查场景合法性、确认初始状态尚未完成，并用参考动作验证任务确实有解；通过后才生成 home_01_laundry_basic.db。",
        3.55,
        4.70,
        8.35,
        0.60,
        size=12,
        color=MUTED,
    )
    add_text(slide, "以上为项目编译产物的真实 SQLite 查询结果。", 0.88, 6.20, 6.2, 0.23, size=8.8, color=FAINT)

    # 09 trajectory
    slide = prs.slides.add_slide(blank)
    add_header(slide, 9, "CASE · 03 / 04", "真实 trajectory：模型如何失败，又如何调整", "middle_output/trajectory.json", VIOLET)
    add_text(slide, "Iteration 4：路线选择失败", 0.85, 1.30, 5.65, 0.32, size=12.5, bold=True)
    add_code(
        slide,
        '{\n  "iteration": 4,\n  "reasoning_content":\n    "现在两件都拿到了。去阳台。",\n  "tool_calls": [{\n    "name": "move_to",\n    "arguments": {"area_id": "balcony"}\n  }],\n  "tool_response": {\n    "status": "failure",\n    "content": "Cannot move from\n      bedroom to balcony — not passable"\n  },\n  "is_final": false\n}',
        0.85,
        1.72,
        5.70,
        4.25,
        name="失败轨迹字段",
        size=9.4,
    )
    add_text(slide, "Iteration 5：读取反馈后改路", 6.78, 1.30, 5.65, 0.32, size=12.5, bold=True)
    add_code(
        slide,
        '{\n  "iteration": 5,\n  "reasoning_content":\n    "不能直接从卧室到阳台。\n     需要先回客厅，再去阳台。",\n  "tool_calls": [{\n    "name": "move_to",\n    "arguments": {"area_id": "living_room"}\n  }],\n  "tool_response": {\n    "status": "success",\n    "content": "Location: 客厅;\n      Holding: 衬衫 + 裤子"\n  },\n  "is_final": false\n}',
        6.78,
        1.72,
        5.70,
        4.25,
        name="恢复轨迹字段",
        size=9.4,
    )
    labels = [
        ("iteration", "动作发生在第几轮", CYAN),
        ("reasoning_content", "模型为何选择该动作", TEAL),
        ("tool_calls / response", "做了什么、如何反馈", VIOLET),
        ("is_final", "这一轮是否结束任务", AMBER),
    ]
    for idx, (title, desc, color) in enumerate(labels):
        x = 0.85 + idx * 2.96
        add_text(slide, title, x, 6.14, 2.70, 0.22, size=9.2, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x, 6.38, 2.70, 0.20, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)

    # 10 最终 DB
    slide = prs.slides.add_slide(blank)
    add_header(slide, 10, "CASE · 04 / 04", "最终不听模型“说完成”，而是查询世界状态", "goal_satisfied = true", AMBER)
    final_data = [
        ["id", "area_id", "container_id", "on_id", "open_state", "device_state"],
        ["clothes_01", "balcony", "washing_machine_01", "NULL", "NULL", "NULL"],
        ["clothes_02", "balcony", "washing_machine_01", "NULL", "NULL", "NULL"],
        ["washing_machine_01", "balcony", "NULL", "NULL", "closed", "running"],
    ]
    add_table(slide, 4, 6, 0.85, 1.43, 11.65, 1.68, final_data, widths=[2.15, 1.45, 2.75, 1.45, 1.85, 2.00], font_size=10.2)
    add_card(slide, 0.85, 3.56, 5.67, 2.05)
    add_text(slide, "任务 JSON 中的三个目标", 1.10, 3.80, 4.95, 0.30, size=13, bold=True)
    goals = [
        "clothes_01.container_id = washing_machine_01",
        "clothes_02.container_id = washing_machine_01",
        "washing_machine_01.device_state = running",
    ]
    for idx, goal in enumerate(goals):
        add_badge(slide, "满足", 1.12, 4.28 + idx * 0.39, 0.62, color=TEAL)
        add_text(slide, goal, 1.88, 4.29 + idx * 0.39, 4.25, 0.28, size=9.8, color=MUTED)
    add_card(slide, 6.78, 3.56, 5.72, 2.05)
    add_text(slide, "程序给出的结论", 7.03, 3.80, 4.95, 0.30, size=13, bold=True)
    add_rich_text(
        slide,
        [("三个条件同时成立，任务结果为 ", MUTED, False), ("complete", AMBER, True), ("。", MUTED, False)],
        7.03,
        4.33,
        4.95,
        0.34,
        size=13,
    )
    add_text(slide, "失败步骤与完整轨迹同时保留，用于后续分析。", 7.03, 4.79, 4.95, 0.35, size=11, color=MUTED)
    add_text(slide, "最终状态来自对 trajectory.json 的真实重放；判定使用项目 goal_dsl.py。", 0.88, 6.10, 8.0, 0.23, size=8.8, color=FAINT)

    # 11 Pelican 闭环
    slide = prs.slides.add_slide(blank)
    add_header(slide, 11, "PELICAN-VL · 01 / 02", "Pelican-VL 是大脑，OpenWorldSandbox 是交互环境", "Brain × Environment", CYAN)
    add_card(slide, 0.58, 1.55, 2.05, 0.88)
    add_text(slide, "任务指令", 0.75, 1.75, 1.71, 0.24, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "+ 当前场景观察", 0.75, 2.06, 1.71, 0.20, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.55, 2.20, 3.40, 2.66, color=FAINT)
    add_loop_node(slide, "Pelican-VL", "理解场景 · 规划任务\n根据反馈调整下一步", 3.35, 2.05, 2.75, TEAL)
    add_loop_node(slide, "OpenWorldSandbox", "执行 17 个语义动作\n检查规则 · 更新 SQLite", 7.65, 2.05, 2.75, VIOLET)
    add_arrow(slide, 5.93, 2.52, 7.58, 2.52, color=CYAN, width=2.3)
    add_text(slide, "高层动作", 6.25, 2.17, 1.02, 0.25, size=10.5, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "move_to / pick_object / ...", 6.02, 2.47, 1.50, 0.20, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 7.68, 4.58, 6.00, 4.58, color=VIOLET, width=2.3)
    add_text(slide, "新观察 / 失败原因", 6.10, 4.79, 1.48, 0.24, size=10.5, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 10.25, 4.33, 11.15, 5.12, color=FAINT)
    add_card(slide, 10.63, 5.12, 2.15, 0.84)
    add_text(slide, "运行产物", 10.82, 5.30, 1.77, 0.23, size=11.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "trajectory · final.db · 结果", 10.82, 5.61, 1.77, 0.20, size=8.7, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "当前主要提供结构化文字观察，测试高层具身推理，而不是完整视觉感知。", 0.67, 6.35, 8.0, 0.22, size=8.8, color=FAINT)
    add_text(slide, "Pelican-VL 1.0：7B–72B 开源具身多模态大脑模型。来源：arXiv:2511.00108。", 0.67, 6.62, 9.5, 0.22, size=8.8, color=FAINT)

    # 12 RL 环境
    slide = prs.slides.add_slide(blank)
    add_header(slide, 12, "PELICAN-VL · 02 / 02", "目标：沉降为 Pelican-VL 的 Agentic RL 执行环境", "Observation → Action → Feedback", AMBER)
    add_card(slide, 0.70, 1.38, 5.88, 4.92)
    add_text(slide, "当前已经具备的交互要素", 1.02, 1.65, 5.25, 0.34, size=14, bold=True)
    existing = [
        "State：SQLite 保存完整符号世界状态",
        "Observation：根据位置和容器状态生成局部观察",
        "Action：17 个高层语义动作",
        "Transition：规则检查后更新世界状态",
        "Termination：完成、放弃、无法继续或超过步数",
        "Trajectory：保存每一步动作和环境反馈",
    ]
    for idx, item in enumerate(existing):
        add_bullet(slide, item, 1.03, 2.18 + idx * 0.53, 5.10, color=TEAL, size=10.5)
    add_card(slide, 6.82, 1.38, 5.80, 4.92)
    add_text(slide, "成为 RL 执行环境需要补齐", 7.14, 1.65, 5.15, 0.34, size=14, bold=True)
    gaps = [
        "统一接口：封装为 reset() / step()",
        "奖励信号：目标、子目标、效率和失败行为",
        "并行运行：大量互不影响的任务环境",
        "模型适配：稳定转换 Pelican-VL 输出",
        "版本与复现：固定环境、任务和奖励版本",
    ]
    for idx, item in enumerate(gaps):
        add_bullet(slide, item, 7.15, 2.18 + idx * 0.53, 5.00, color=AMBER, size=10.5, symbol="+")
    add_card(slide, 7.15, 5.05, 5.05, 0.82, fill=CARD_2, line=AMBER)
    add_text(slide, "目标形态", 7.38, 5.21, 0.92, 0.24, size=11, color=AMBER, bold=True)
    add_text(slide, "Sandbox 负责状态变化、奖励计算和轨迹记录。", 8.25, 5.15, 3.68, 0.38, size=10.5, color=MUTED)
    add_summary_bar(slide, "Pelican-VL 持续观察、行动并接收反馈；环境提供稳定、可审计、可复现的训练闭环。", 6.48, color=AMBER)

    # 13 建设计划
    slide = prs.slides.add_slide(blank)
    add_header(slide, 13, "建设计划", "围绕三条主线，把最小闭环沉降成长期能力", "Environment × Tasks × Agentic RL", AMBER)
    plans = [
        ("01", "环境可靠性", CYAN, "补齐核心动作、目标判定和编译门禁的回归测试；实现环境事件、稳定重置和版本管理。", "相同任务可恢复到同一初始状态，规则修改不破坏历史任务。"),
        ("02", "任务体系", TEAL, "按单项能力、多步组合、失败恢复、复杂约束组织任务，并加入状态与路线扰动。", "不同难度稳定拉开模型表现，失败可对应到明确能力。"),
        ("03", "Agentic RL 接口", VIOLET, "封装 reset / step，设计可审计奖励，支持并行环境、动作适配和统一轨迹格式。", "持续返回 observation、reward、done 和训练所需轨迹。"),
    ]
    for idx, (number, title, color, desc, acceptance) in enumerate(plans):
        x = 0.68 + idx * 4.18
        add_card(slide, x, 1.48, 3.82, 4.35)
        add_badge(slide, number, x + 0.25, 1.78, 0.62, color=color)
        add_text(slide, title, x + 1.02, 1.75, 2.50, 0.37, size=15, bold=True)
        add_line(slide, x + 0.25, 2.40, x + 3.55, 2.40, color=LINE)
        add_text(slide, desc, x + 0.27, 2.70, 3.28, 1.28, size=11.2, color=MUTED, valign=MSO_ANCHOR.TOP)
        add_card(slide, x + 0.27, 4.47, 3.28, 0.90, fill=CARD_2, line=color)
        add_text(slide, "验收", x + 0.47, 4.61, 0.55, 0.24, size=10.5, color=color, bold=True)
        add_text(slide, acceptance, x + 1.08, 4.55, 2.27, 0.54, size=9.5, color=MUTED)
    add_summary_bar(slide, "预期交付：可重复的具身推理环境、分层任务资产、可解释能力报告，以及可供 Agentic RL 使用的执行接口。", 6.25, color=AMBER)

    # 14 Smoke test 与 baseline
    slide = prs.slides.add_slide(blank)
    add_header(slide, 14, "验证结果 · 01 / 03", "Smoke test：闭环已跑通，正式榜单与探索性测试分开", "GLM-5.1 baseline × DeepSeek-V4-Pro smoke", TEAL)
    add_card(slide, 0.72, 1.42, 5.86, 4.62)
    add_badge(slide, "榜单", 1.02, 1.74, 0.72, color=TEAL)
    add_text(slide, "GLM-5.1 · 标准 baseline", 1.92, 1.70, 3.95, 0.34, size=15, bold=True)
    add_text(slide, "7 个任务 × 每任务 3 次 = 21 runs", 1.03, 2.24, 5.05, 0.26, size=10.8, color=MUTED)
    glm_metrics = [
        ("Success", "71.4%", TEAL),
        ("pass@3", "100%", CYAN),
        ("SPL", "0.581", VIOLET),
        ("Step Ratio", "1.45", AMBER),
    ]
    for idx, (label, value, color) in enumerate(glm_metrics):
        x = 1.02 + (idx % 2) * 2.58
        y = 2.72 + (idx // 2) * 1.02
        add_card(slide, x, y, 2.32, 0.78, fill=CARD_2, line=color)
        add_text(slide, label, x + 0.18, y + 0.13, 1.15, 0.20, size=9.5, color=MUTED)
        add_text(slide, value, x + 1.15, y + 0.11, 0.90, 0.28, size=15, color=color, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "榜单口径：同一协议批量运行，由 verify/report 自动汇总。", 1.04, 5.48, 4.95, 0.28, size=9.5, color=FAINT)

    add_card(slide, 6.78, 1.42, 5.83, 4.62)
    add_badge(slide, "探索", 7.08, 1.74, 0.72, color=VIOLET)
    add_text(slide, "DeepSeek-V4-Pro · 功能 smoke", 7.98, 1.70, 4.00, 0.34, size=15, bold=True)
    ds_rows = [
        ("洗衣任务", "complete", "12 / 40", TEAL),
        ("厨房清洁", "complete", "22 / 30", TEAL),
        ("厨房安全复原", "exceeded", "40 / 40", RED),
    ]
    for idx, (task_name, result, steps, color) in enumerate(ds_rows):
        y = 2.54 + idx * 0.72
        add_text(slide, task_name, 7.10, y, 2.20, 0.28, size=10.5, color=INK, bold=True)
        add_badge(slide, result, 9.38, y - 0.01, 1.12, color=color)
        add_text(slide, steps, 10.82, y, 1.18, 0.28, size=10.5, color=MUTED, align=PP_ALIGN.RIGHT)
    add_card(slide, 7.08, 4.84, 5.10, 0.76, fill=CARD_2, line=VIOLET)
    add_text(slide, "未纳入 leaderboard", 7.32, 4.98, 1.72, 0.24, size=11, color=VIOLET, bold=True)
    add_text(slide, "单次探索、任务子集与裁判设置不同，不与 GLM 横向比较。", 9.05, 4.91, 2.86, 0.39, size=9.2, color=MUTED)
    add_summary_bar(slide, "结论：harness 已覆盖生成动作、环境反馈、终态判分与指标汇总；跨模型结论仍需统一协议复测。", 6.45, color=TEAL)
    add_text(slide, "来源：outputs/runs/baseline/leaderboard.json；middle_output/* DeepSeek-V4-Pro 测试记录。", 0.78, 6.91, 9.5, 0.18, size=8.3, color=FAINT)

    # 15 数据合成压测
    slide = prs.slides.add_slide(blank)
    add_header(slide, 15, "验证结果 · 02 / 03", "可规模化合成：并发 4 接近线性，并发 8 开始饱和", "真实 API 压测 · 67 候选 → 34 合格任务", CYAN)
    synth_table = [
        ["并发", "候选", "合格", "合格任务/min", "扩展效率", "Compile 通过率", "成本/合格任务"],
        ["1", "24", "13", "1.78", "100.0%", "54.2%", "$0.0288"],
        ["4", "20", "12", "6.51", "91.5%", "60.0%", "$0.0265"],
        ["8", "23", "9", "6.81", "47.9%", "39.1%", "$0.0437"],
    ]
    add_table(
        slide,
        4,
        7,
        0.72,
        1.38,
        11.89,
        2.18,
        synth_table,
        widths=[1.15, 1.35, 1.35, 2.05, 1.65, 2.10, 2.24],
        font_size=10.0,
    )
    add_text(slide, "合格任务吞吐（tasks / min）", 0.82, 3.92, 4.20, 0.28, size=12.2, bold=True)
    rates = [(1, 1.7784, TEAL), (4, 6.5055, CYAN), (8, 6.8129, VIOLET)]
    for idx, (concurrency, rate, color) in enumerate(rates):
        y = 4.42 + idx * 0.57
        add_text(slide, f"并发 {concurrency}", 0.84, y, 0.84, 0.24, size=9.5, color=MUTED)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.76, y + 0.03, rate * 0.58, 0.19, fill=color, line=color)
        add_text(slide, f"{rate:.2f}", 5.88, y, 0.58, 0.24, size=10, color=color, bold=True, align=PP_ALIGN.RIGHT)
    add_card(slide, 6.78, 3.90, 5.83, 1.88)
    add_text(slide, "这次压测真正证明了什么", 7.08, 4.16, 4.96, 0.30, size=13.5, bold=True)
    add_bullet(slide, "并发 1 → 4：吞吐提升 3.66×，扩展效率 91.5%", 7.08, 4.66, 4.95, color=CYAN, size=10.2)
    add_bullet(slide, "并发 4 → 8：吞吐仅提升 4.7%，API/worker 已出现饱和", 7.08, 5.09, 4.95, color=AMBER, size=10.2)
    add_summary_bar(slide, "14 个场景、34 个 compile 合格任务已进入 staging；每条请求、失败与编译结果均可由 events.jsonl 回溯。", 6.30, color=CYAN)
    add_text(slide, "来源：outputs/reports/synth_report_20260825_132735.json。小样本压测，不等同于长期稳定产能。", 0.78, 6.91, 10.0, 0.18, size=8.3, color=FAINT)

    # 16 实测、外推与边界
    slide = prs.slides.add_slide(blank)
    add_header(slide, 16, "验证结果 · 03 / 03", "Scalable 潜力：把实测、外推和理论上限分开讲", "Measured ≠ Extrapolated ≠ Theoretical", AMBER)
    capacity_cards = [
        ("实测", "408.77", "合格任务 / 小时", TEAL, "并发 8 的短时压测吞吐\n6.81 tasks/min × 60"),
        ("基于实测外推", "≈ 6,540", "合格任务 / 天", CYAN, "假设每日运行 16 小时\n且吞吐与通过率不衰减"),
        ("理论上限", "1,044.66", "合格任务 / 小时", VIOLET, "未验证：100% 过门禁\n无限 API 配额、无重试"),
    ]
    for idx, (label, value, unit, color, note) in enumerate(capacity_cards):
        x = 0.68 + idx * 4.18
        add_card(slide, x, 1.40, 3.82, 2.52)
        add_badge(slide, label, x + 0.26, 1.72, 1.28, color=color)
        add_text(slide, value, x + 0.26, 2.16, 3.28, 0.55, size=25, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, unit, x + 0.26, 2.72, 3.28, 0.26, size=10.2, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, note, x + 0.30, 3.13, 3.20, 0.58, size=9.2, color=FAINT, align=PP_ALIGN.CENTER)
    add_card(slide, 0.68, 4.28, 7.79, 1.62)
    add_text(slide, "扩产闭环已经成形", 0.98, 4.55, 2.05, 0.30, size=13.5, bold=True)
    flow = [
        ("LLM API", CYAN),
        ("JSON 候选", TEAL),
        ("Compile 门禁", VIOLET),
        ("Staging", AMBER),
        ("人工抽检", RED),
    ]
    for idx, (name, color) in enumerate(flow):
        x = 0.96 + idx * 1.40
        add_badge(slide, name, x, 5.11, 1.08, color=color)
        if idx < len(flow) - 1:
            add_arrow(slide, x + 1.10, 5.25, x + 1.34, 5.25, color=FAINT)
    add_card(slide, 8.70, 4.28, 3.91, 1.62, fill=CARD_2, line=AMBER)
    add_text(slide, "当前瓶颈 / 下一步", 8.98, 4.55, 3.30, 0.30, size=13.5, color=AMBER, bold=True)
    add_text(slide, "扩大样本验证稳定性；修复 3 次 worker crash；测 API 限流与人工审核成本。", 8.98, 5.04, 3.25, 0.60, size=10.2, color=MUTED, valign=MSO_ANCHOR.TOP)
    add_summary_bar(slide, "对领导的承诺口径：已证明自动扩产链路可行；日产 6,540 是小样本线性外推，不是稳定 SLA。", 6.30, color=AMBER)
    add_text(slide, "压测限制：每档仅 20–24 个候选；并发 1 基线受服务端降速影响；并发 8 有 2 次 worker crash。", 0.78, 6.91, 10.8, 0.18, size=8.3, color=FAINT)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"已生成：{OUTPUT}")
    print(f"页数：{len(prs.slides)}")
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"第 {idx:02d} 页：{len(slide.shapes)} 个可编辑对象")


if __name__ == "__main__":
    build()
